#pip install google-api-python-client python--pptx langchain langchain_google_genai PyMuPDF PyPDF2 google-search-results --user

import streamlit as st
from pptx import Presentation
from langchain import PromptTemplate
from langchain_google_genai import ChatGoogleGenerativeAI
from googleapiclient.discovery import build
from serpapi import GoogleSearch
import time
import PyPDF2

# Streamlit App Title
st.title("AI Tutor")

# Function to read API keys from config file
def load_api_keys(config_file='config.txt'):
    api_keys = {}
    with open(config_file, 'r') as f:
        for line in f:
            if '=' in line:
                key, value = line.strip().split('=')
                api_keys[key] = value
    return api_keys

# Load API keys from config.txt
api_keys = load_api_keys()

# Use the keys in your application
gemini_api_key = api_keys['GEMINI_API_KEY']
youtube_api_key = api_keys['YOUTUBE_API_KEY']
serpapi_api_key = api_keys['SERPAPI_API_KEY']

# Extract text from PPTX file
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Extract text from PDF file
def extract_text_from_pdf(file_path):
    reader = PyPDF2.PdfReader(file_path)
    text = ""
    for page_num in range(len(reader.pages)):
        page = reader.pages[page_num]
        text += page.extract_text()
    return text

# Create prompt template for Gemini
llm_prompt_template = """Extract the most important keywords from the following text:
"{text}"
KEYWORDS:"""
llm_prompt = PromptTemplate.from_template(llm_prompt_template)

# Summarize text with Gemini and extract keywords
def get_keywords_with_gemini(text):
    llm = ChatGoogleGenerativeAI(
        model="gemini-pro",
        temperature=0.7,
        top_p=0.85,
        google_api_key=gemini_api_key
    )
    prompt = llm_prompt.format(text=text)
    response = llm.invoke(prompt)
    keywords = response.content.split(", ")  # Assuming Gemini returns comma-separated keywords
    return keywords[:5]  # Limit to top 5 keywords

# Search YouTube for related videos with retries
def search_youtube(query, retries=3, wait_time=2):
    youtube = build('youtube', 'v3', developerKey=youtube_api_key)
    
    for attempt in range(retries):
        #st.write(f"Query used for YouTube search: {query}")  # Log the search query

        request = youtube.search().list(
            q=query,
            part='snippet',
            type='video',
            order='viewCount',  # Order by view count
            maxResults=10
        )
        response = request.execute()

        video_ids = [item['id']['videoId'] for item in response['items']]

        if video_ids:  # Check if the result is not empty
            # Fetch video statistics
            videos = []
            stats_request = youtube.videos().list(
                part='snippet,statistics',
                id=','.join(video_ids)
            )
            stats_response = stats_request.execute()

            for item in stats_response['items']:
                video_title = item['snippet']['title']
                video_id = item['id']
                video_url = f"https://www.youtube.com/watch?v={video_id}"
                video_views = item['statistics'].get('viewCount', 'N/A')
                videos.append((video_title, video_url, int(video_views) if video_views != 'N/A' else 0))

            # Sort by views if available
            videos.sort(key=lambda x: x[2], reverse=True)
            return videos[:5]  # Return only the top 5 results
        else:
            st.warning(f"No YouTube results found, retrying... (Retry {attempt+1}/{retries})")
            time.sleep(wait_time)  # Wait before retrying

    st.error("YouTube search failed after multiple attempts.")
    return []  # Return empty if no results after retries

# Search web for related blogs/websites
def search_web(query):
    #st.write(f"Query used for web search: {query}")  # Log the search query
    search = GoogleSearch({
        'q': query,
        'num': 10,
        'api_key': serpapi_api_key,
        'gl': 'us',
        'hl': 'en'
    })
    results = search.get_dict()

    # Log the raw search response for debugging
    #st.write(f"Raw web search response: {results}")

    blogs = []
    for result in results.get('organic_results', []):
        if 'link' in result and 'title' in result:
            title = result['title']
            link = result['link']
            snippet = result.get('snippet', '')
            blogs.append((title, link, snippet))

    return blogs[:5]  # Limit to top 5 results

# File uploader in Streamlit (automatically detects file type)
uploaded_file = st.file_uploader("Upload your file (PPTX or PDF)", type=['pptx', 'pdf'])

if uploaded_file is not None:
    with st.spinner("Extracting text and processing..."):
        try:
            # Automatically detect the file type and extract text accordingly
            file_extension = uploaded_file.name.split('.')[-1].lower()

            if file_extension == 'pptx':
                text = extract_text_from_pptx(uploaded_file)
            elif file_extension == 'pdf':
                text = extract_text_from_pdf(uploaded_file)
            else:
                st.error("Unsupported file format. Please upload a PPTX or PDF file.")
                st.stop()

            # Get keywords using Gemini
            keywords = get_keywords_with_gemini(text)
            query = " ".join(keywords)

            # Sidebar to choose between YouTube or Blogs/Website results
            choice = st.sidebar.radio("", ["YouTube Results", "Blogs/Website Results"])

            if choice == "YouTube Results":
                st.subheader("YouTube Search Results")
                youtube_results = search_youtube(query)
                if youtube_results:
                    for title, url, views in youtube_results:
                        st.markdown(f"**{title}**<br>{url}<br>Views: {views}<br>", unsafe_allow_html=True)
                else:
                    st.write("No YouTube results found.")

            elif choice == "Blogs/Website Results":
                st.subheader("Web Search Results")
                web_results = search_web(query)
                if web_results:
                    for title, link, snippet in web_results:
                        st.markdown(f"**{title}**<br>{link}<br>{snippet}<br>", unsafe_allow_html=True)
                else:
                    st.write("No web search results found.")

        except Exception as e:
            st.error(f"An error occurred: {e}")
else:
    st.info("Please upload a PPTX or PDF file to proceed.")
